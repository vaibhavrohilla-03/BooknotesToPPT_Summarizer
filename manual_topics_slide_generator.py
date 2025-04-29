# Save as manual_topics_slide_generator.py
import os
import re
import numpy as np
from PyPDF2 import PdfReader
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import torch
from pptx import Presentation
from pptx.util import Inches, Pt
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.probability import FreqDist
import fitz 


try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)

class ManualTopicSlideGenerator:
    def __init__(self, textbook_path, topics=None):
        self.textbook_path = textbook_path
        self.syllabus_topics = topics or []
        self.textbook_content = {}
        self.topic_matches = {}
        self.presentation = None
        self.images = {}
    
    def set_topics(self, topics):
       
        if isinstance(topics, str):
           
            self.syllabus_topics = [topic.strip() for topic in topics.split('\n') if topic.strip()]
        elif isinstance(topics, list):
            
            self.syllabus_topics = topics
        else:
            raise ValueError("Topics must be a list of strings or a newline-separated string")
        
        print(f"Set {len(self.syllabus_topics)} manual topics:")
        for topic in self.syllabus_topics[:5]:
            print(f"- {topic}")
        if len(self.syllabus_topics) > 5:
            print(f"- ... and {len(self.syllabus_topics)-5} more")
        
        return self.syllabus_topics
    
    def extract_textbook_content(self):
       
        print("Extracting content from textbook...")
        doc = fitz.open(self.textbook_path)
        
     
        full_text = ""
        current_section = "Introduction"
        self.textbook_content[current_section] = ""
        
        for page_num in range(len(doc)):
            print(f"Processing page {page_num+1}/{len(doc)}...")
            page = doc.load_page(page_num)
            blocks = page.get_text("dict")["blocks"]
            
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span["text"]
                            font_size = span["size"]
                            is_bold = span["flags"] & 2 > 0  
                            
                         
                            if (font_size > 11 and len(text) < 100) or is_bold:
                                potential_heading = text.strip()
                           
                                if len(potential_heading) > 3:
                                    current_section = potential_heading
                                    if current_section not in self.textbook_content:
                                        self.textbook_content[current_section] = ""
                            else:
                                self.textbook_content[current_section] += text + " "
        
     
        self.images = {}
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                  
                    nearest_heading = list(self.textbook_content.keys())[-1]  
                    if nearest_heading not in self.images:
                        self.images[nearest_heading] = []
                    
                  
                    img_filename = f"temp_img_{page_num}_{img_index}.png"
                    with open(img_filename, "wb") as img_file:
                        img_file.write(image_bytes)
                    self.images[nearest_heading].append(img_filename)
                except Exception as e:
                    print(f"Error extracting image: {str(e)}")
        
        
        for section in list(self.textbook_content.keys()):
            if len(self.textbook_content[section].split()) < 20:
                del self.textbook_content[section]
        
        print(f"Extracted {len(self.textbook_content)} sections from textbook")
        return self.textbook_content
    
    def match_topics_to_content(self):
        """Use NLP to match syllabus topics with textbook sections"""
        print("Matching topics to textbook sections...")
       
        if not self.syllabus_topics:
            print("No topics provided. Please set topics using set_topics() method.")
            return {}
            
        if not self.textbook_content:
            print("No textbook content extracted.")
            return {}
        
       
        headings = list(self.textbook_content.keys())
        
        # Use TF-IDF for matching
        vectorizer = TfidfVectorizer(stop_words='english')
        
       
        heading_samples = []
        for heading in headings:
            content = self.textbook_content[heading]
            
            sample = heading + " " + " ".join(content.split()[:200])
            heading_samples.append(sample)
        
       
        try:
            corpus = heading_samples + self.syllabus_topics
            tfidf_matrix = vectorizer.fit_transform(corpus)
            
            
            heading_tfidf = tfidf_matrix[:len(headings)]
            topic_tfidf = tfidf_matrix[len(headings):]
            
            
            similarity_matrix = cosine_similarity(topic_tfidf, heading_tfidf)
            
          
            for i, topic in enumerate(self.syllabus_topics):
               
                scores = similarity_matrix[i]
               
                top_indices = np.argsort(scores)[::-1]
                
               
                matches = []
                for idx in top_indices:
                    if scores[idx] > 0.05 and len(matches) < 3:  
                        matches.append(headings[idx])
                
                self.topic_matches[topic] = matches if matches else [headings[0]]  
                
                # Debug output
                print(f"Topic '{topic}' matched to sections: {matches}")
        
        except Exception as e:
            print(f"Error during topic matching: {str(e)}")
           
            for topic in self.syllabus_topics:
                self.topic_matches[topic] = headings[:3] 
        
        return self.topic_matches
    
    def summarize_content(self, section_text, max_sentences=5):
    
        sentences = sent_tokenize(section_text)
        
        if len(sentences) <= max_sentences:
            return section_text
        
      
        stop_words = set(stopwords.words('english'))
        words = word_tokenize(section_text.lower())
        important_words = [word for word in words if word.isalnum() and word not in stop_words]
        word_freq = FreqDist(important_words)
        
     
        sentence_scores = []
        for i, sent in enumerate(sentences):
       
            position_score = 1.0 if i < len(sentences) // 3 else 0.5
            
          
            words_in_sentence = [word.lower() for word in word_tokenize(sent) if word.lower() in word_freq]
            freq_score = sum([word_freq[word] for word in words_in_sentence]) / max(1, len(words_in_sentence))
            
         
            length = len(word_tokenize(sent))
            length_score = 1.0 if 5 <= length <= 25 else 0.5
            
          
            total_score = (freq_score * 0.6) + (position_score * 0.3) + (length_score * 0.1)
            sentence_scores.append((total_score, i, sent))
        
       
        top_sentences = sorted(sentence_scores, key=lambda x: x[0], reverse=True)[:max_sentences]
        top_sentences = sorted(top_sentences, key=lambda x: x[1])
        
      
        summary = ' '.join([sent for _, _, sent in top_sentences])
        return summary
    
    def extract_definitions_and_formulas(self, text):
      
        definitions = re.findall(r'([A-Z][^.!?]*(?:is defined as|refers to|means|is called)[^.!?]*\.)', text)
       
        formulas = re.findall(r'([^.!?]*(?:\$.*\$|\d+\s*(?:=|≡|≈|≠|<|>|≤|≥)[^.!?]*\))[^.!?]*\.)', text)
        
       
        sentences = sent_tokenize(text)
        definition_keywords = ['defined', 'known as', 'refers to', 'called', 'means', 'denotes']
        additional_definitions = []
        
        for sent in sentences:
            if any(keyword in sent.lower() for keyword in definition_keywords):
                if sent not in definitions:
                    additional_definitions.append(sent)
        
        combined_definitions = definitions + additional_definitions
        return {
            'definitions': combined_definitions[:3],  
            'formulas': formulas[:3]
        }
    
    def extract_key_points(self, text, max_points=5):
       
        sentences = sent_tokenize(text)
        
        if not sentences:
            return []
        
      
        if len(sentences) <= 2:
            return sentences
        
       
        key_indicators = [
            "important", "key", "critical", "essential", "fundamental", 
            "significant", "crucial", "primary", "main", "central",
            "note that", "remember", "notably"
        ]
        
        potential_key_points = []
        
     
        for sent in sentences:
            sent = sent.strip()
            lower_sent = sent.lower()
            
          
            if len(sent.split()) < 4 or len(sent.split()) > 40:
                continue
                
          
            if any(indicator in lower_sent for indicator in key_indicators):
                potential_key_points.append(sent)
                
        
            if " is " in lower_sent and len(sent.split()) < 25:
                potential_key_points.append(sent)
                
           
            if re.search(r'\d+', sent) and len(sent.split()) < 30:
                potential_key_points.append(sent)
        
    
        if len(potential_key_points) < max_points:
       
            if sentences[0] not in potential_key_points:
                potential_key_points.append(sentences[0])
                
          
            positions = [
                int(len(sentences) * 0.25),
                int(len(sentences) * 0.5),
                int(len(sentences) * 0.75)
            ]
            
            for pos in positions:
                if len(potential_key_points) < max_points and sentences[pos] not in potential_key_points:
                    potential_key_points.append(sentences[pos])
        
        # Limit to max_points
        return potential_key_points[:max_points]
    
    def generate_slides(self, output_path):
        """Generate PowerPoint slides based on matched content"""
        print("Generating PowerPoint slides...")
        self.presentation = Presentation()
        
       
        title_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Generated Course Slides"
        subtitle.text = "Based on syllabus topics and textbook content"
        
      
        toc_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        toc_title = toc_slide.shapes.title
        toc_content = toc_slide.placeholders[1]
        toc_title.text = "Table of Contents"
        toc_text = ""
        for i, topic in enumerate(self.syllabus_topics):
            toc_text += f"{i+1}. {topic}\n"
        toc_content.text = toc_text
        
   
        if not self.topic_matches:
           
            error_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
            error_title = error_slide.shapes.title
            error_content = error_slide.placeholders[1]
            error_title.text = "Error: No Content Matched"
            error_content.text = (
                "The system could not match syllabus topics to textbook content.\n\n"
                "Possible reasons:\n"
                "• PDF format not properly recognized\n"
                "• No clearly defined sections in textbook\n"
                "• Topics don't match textbook terminology\n\n"
                "Try reviewing the extracted content and topic matches in console output."
            )
            self.presentation.save(output_path)
            return output_path
        
    
        for topic_idx, topic in enumerate(self.syllabus_topics):
          
            topic_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
            title = topic_slide.shapes.title
            title.text = f"Topic {topic_idx+1}: {topic}"
            
         
            matched_sections = self.topic_matches.get(topic, [])
            if not matched_sections:
               
                content = topic_slide.placeholders[1]
                content.text = "No matching content found in textbook for this topic."
                continue
            
           
            overview = topic_slide.placeholders[1]
            overview_text = "Related sections in textbook:\n"
            for section in matched_sections:
               
                content_words = len(self.textbook_content.get(section, "").split())
                overview_text += f"• {section} ({content_words} words)\n"
            overview.text = overview_text
            
        
            for section_idx, section in enumerate(matched_sections):
                content = self.textbook_content.get(section, "")
                if not content or len(content.split()) < 20: 
                    continue
                
              
                content_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[2])
                title = content_slide.shapes.title
                content_placeholder = content_slide.placeholders[1]
                
                title.text = f"{topic} - {section}"
                
             
                key_points = self.extract_key_points(content, max_points=5)
                definitions = self.extract_definitions_and_formulas(content)
                
               
                slide_text = ""
              
                if key_points:
                    slide_text += "Key Points:\n"
                    for point in key_points:
                        slide_text += f"• {point}\n"
                    slide_text += "\n"
                
              
                if definitions['definitions']:
                    slide_text += "Key Definitions:\n"
                    for defn in definitions['definitions'][:2]:  
                        slide_text += f"• {defn}\n"
                    slide_text += "\n"
                
               
                if definitions['formulas']:
                    slide_text += "Key Formulas:\n"
                    for formula in definitions['formulas'][:2]:  
                        slide_text += f"• {formula}\n"
             
                if not slide_text.strip():
                  
                    sentences = sent_tokenize(content)
                    sample_text = " ".join(sentences[:3]) + "..."
                    slide_text = f"Content extract:\n\n{sample_text}"
                
                content_placeholder.text = slide_text
                
               
                if section in self.images and self.images[section]:
                    try:
                        image_slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[3])
                        image_title = image_slide.shapes.title
                        image_title.text = f"{topic} - {section} (Visuals)"
                        
                      
                        for i, img_path in enumerate(self.images[section][:2]):
                            if os.path.exists(img_path):
                                left = Inches(1 if i == 0 else 5)
                                top = Inches(2)
                                image_slide.shapes.add_picture(img_path, left, top, width=Inches(4))
                    except Exception as e:
                        print(f"Error adding image slide: {str(e)}")
      
        try:
            self.presentation.save(output_path)
            print(f"Slides successfully saved to {output_path}")
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")
            
            alt_path = "backup_slides.pptx"
            try:
                self.presentation.save(alt_path)
                print(f"Slides saved to alternate location: {alt_path}")
                return alt_path
            except:
                print("Failed to save presentation at all.")
        
       
        for section, image_list in self.images.items():
            for img_path in image_list:
                if os.path.exists(img_path):
                    try:
                        os.remove(img_path)
                    except:
                        pass
        
        return output_path
    
    def run_pipeline(self, output_path="generated_slides.pptx"):
        
        try:
            print("\nStarting slide generation pipeline...")
            print(f"Textbook: {self.textbook_path}")
            
           
            if not os.path.exists(self.textbook_path):
                print(f"ERROR: Textbook file not found at {self.textbook_path}")
                return None
            
            
            if not self.syllabus_topics:
                print("ERROR: No topics provided. Please set topics using set_topics() method.")
                return None
            
            print("\nStep 1: Extracting textbook content...")
            content = self.extract_textbook_content()
            content_sections = len(content)
            print(f"Found {content_sections} sections in textbook")
            
      
            total_words = sum(len(text.split()) for text in content.values())
            print(f"Total extracted words: {total_words}")
            
            if content_sections < 2 or total_words < 100:
                print("WARNING: Very little content extracted from textbook.")
                print("This may result in empty or low-quality slides.")
            
            print("\nStep 2: Matching topics to content...")
            matches = self.match_topics_to_content()
            total_matches = sum(len(sections) for sections in matches.values())
            print(f"Made {total_matches} topic-section matches")
            
            print("\nStep 3: Generating slides...")
            output_file = self.generate_slides(output_path)
            
            print(f"\n✓ Slides generated successfully at {output_file}")
            return output_file
                
        except Exception as e:
            print(f"\nError in pipeline: {str(e)}")
            import traceback
            traceback.print_exc()
            return None
    
    @staticmethod
    def generate(textbook_path, topics, output_path="generated_slides.pptx"):
        
        generator = ManualTopicSlideGenerator(textbook_path)
        generator.set_topics(topics)
        return generator.run_pipeline(output_path)